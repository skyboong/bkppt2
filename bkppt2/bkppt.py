print(">>> 0.0.1")

import re 
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor  # 색상 설정
from pptx.util import Mm

from pathlib import Path
from functools import partial

from icecream import ic 

#import re 
#from datetime import datetime 
from pandas import DataFrame
#from pathlib import Path 

SLD_LAYOUT_TITLE_AND_CONTENT = 1

# 색상 딕셔너리로 정리
COLORS = {
    "gray_dark": RGBColor(102, 102, 102),
    "gray_light": RGBColor(200, 200, 200),
    "navy_blue": RGBColor(0, 0, 128),
    "royal_blue": RGBColor(65, 105, 225),
    "crimson_red": RGBColor(220, 20, 60) ,     # Crimson Red (#DC143C)
    "emerald_green" : RGBColor(80, 200, 120) ,  # Emerald Green (#50C878)
    "orange" : RGBColor(255, 165, 0) ,          # Orange (#FFA500)
    "black" : RGBColor(0, 0, 0),
    "gray" : RGBColor(102, 102, 102),
    # 회색 계열 색상 정의
    "gray_light": RGBColor(200, 200, 200)  ,     # 밝은 회색 (Light Gray)
    "gray_medium": RGBColor(150, 150, 150)  ,    # 중간 회색 (Medium Gray)
    "gray_very_dark": RGBColor(51, 51, 51)  ,    # 매우 어두운 회색 (Very Dark Gray)

    # 추가 회색 계열 색상
    "gray_silver": RGBColor(192, 192, 192)  ,    # 실버 (Silver Gray)
    "gray_charcoal": RGBColor(77, 77, 77) ,      # 차콜 그레이 (Charcoal Gray)
    "gray_slate": RGBColor(112, 128, 144)   ,    # 슬레이트 그레이 (Slate Gray)
    "gray_dim": RGBColor(105, 105, 105)   ,
}

COLORS_NAME  = {
    "deepblue": "#1E3A8A",
    "darkgray": "#374151",
    "deeporange": "#D97706",

    "blue": "#0000FF",
    "green": "#008000",
    "red": "#FF0000",
    "cyan": "#00FFFF",
    "magenta": "#FF00FF",
    "yellow": "#FFFF00",
    "black": "#000000",
    "white": "#FFFFFF",
    "gray": "#808080",
    "orange": "#FFA500",
    "purple": "#800080",
    "brown": "#A52A2A",
    "pink": "#FFC0CB",
    "lime": "#00FF00",
    "indigo": "#4B0082",
    "violet": "#8A2BE2",
    "gold": "#FFD700",
    "silver": "#C0C0C0",
    "beige": "#F5F5DC",
    "teal": "#008080",
    "navy": "#000080",
    "turquoise": "#40E0D0",
    "orchid": "#DA70D6",

        # 추가된 엷은 계열 색상들
    "lightblue": "#ADD8E6",
    "lightgreen": "#90EE90",
    "lightpink": "#FFB6C1",
    "lightyellow": "#FFFFE0",
    "lightcyan": "#E0FFFF",
    "lightcoral": "#F08080",
    "lightseagreen": "#20B2AA",
    "lightgoldenrodyellow": "#FAFAD2",
    "lavender": "#E6E6FA",
    "thistle": "#D8BFD8",
    "peachpuff": "#FFDAB9",
    "mistyrose": "#FFE4E1",
    "honeydew": "#F0FFF0",
    "papayawhip": "#FFEFD5",
    "azure": "#F0FFFF",
    "mintcream": "#F5FFFA"
}

# 기본 설정 상수
DEFAULT_TEXTBOX_WIDTH = 12.33
DEFAULT_TEXTBOX_HEIGHT = 6.0
DEFAULT_FONT_NAME = '나눔고딕'
DEFAULT_COLOR_BOLD = '#ff0000'

def create_pppt(prefix:str='ppt_', 
              slides_data:list=[], slide_width:float=13.33,
                slide_height:float=7.5, 
                default_textbox_width=DEFAULT_TEXTBOX_WIDTH, 
                default_textbox_height=DEFAULT_TEXTBOX_HEIGHT,
                font_size_title = 20,
                font_size_caption_title=6,
                font_size_level_0 = 16,
                font_size_level_1 = 10,
                font_size_level_2 = 8,
                font_size_level_3 = 7,
                font_size_level_4 = 5,
                font_size_default = 10,
              dir_figure='images',
              font_name=DEFAULT_FONT_NAME,
              space_before=5,
              space_after=5,
              space_before1=1,
              space_after1=1,
              color_bold=DEFAULT_COLOR_BOLD,
              color_background =  "#FFFFFF", # "#4169E1",
              auto_paragraphs_two = False,
              auto_paragraphs_threshold=500,
              auto_figure_position=False,
              auto_figure_vertical=True,
              auto_figure_xp=10,
              auto_figure_yp=2,
              auto_figure_height=1,
              auto_figure_delta=0.2):
    print(f">>> prefix={prefix}")
    ic()
    ic.disable()
    current_dir = Path(__file__).parent
    # 프레젠테이션 생성
   
    default_textbox_width1 = Inches(default_textbox_width)
    default_textbox_width_half = Inches(default_textbox_width/2)
    default_textbox_height1 = Inches(default_textbox_height)
    ic(default_textbox_height1)
    print("....")

    pos_x = auto_figure_xp
    pos_y = auto_figure_yp

    default_left = Inches(0.5)
    default_top = Inches(1)
 
    level_font_sizes = {
        0: Pt(font_size_level_0),  # 제목 수준
        1: Pt(font_size_level_1),
        2: Pt(font_size_level_2),
        3: Pt(font_size_level_3),
        4: Pt(font_size_level_4)   # 가장 낮은 수준
    }
    
    prs = Presentation()
    prs.slide_width = Inches(slide_width)  # 원하는 가로 폭 (기본 13.33보다 넓게)
    prs.slide_height = Inches(slide_height)

    # 슬라이드 추가 (텍스트 슬라이드)
    #slide_no = 0 
    #for i in range(len(slides_data)):

    i = 0 
    #print(f"i={i}")
    #for slide_data in slides_data[:2]:
    for slide_data in slides_data:
        i += 1
        #print(f"ith={i}")
        #layout_idx = find_slide_layout_number(slides_data[i]['body'],3)
        layout_idx = extract_slide_layout(slide_data['body'], default_index=3)
        #print(f"* slide layout index= {layout_idx}")
        
        # (1) 슬라이드 추가 
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = hex_to_rgb_color(color_background)


        title = slide.shapes.title
        body_text = clean_text(text=slide_data['body'], kind='layout_index')
        body_text_list = [each for each in body_text.split('\n') ]
   
        body_11 = body_text 
        body_11 = remove_ptn_in_text(r"^<section:", body_11)
        body_11 = remove_ptn_in_text(r"^<width:", body_11)
        body_11 = remove_ptn_in_text(r"^<fig:", body_11)
        
        # @@@을 직접입력해 둔 경우는 직접입력한 것 우선으로 처리해 줌 
        n_tf_at = 0 
        for each_body_text in body_text.split('\n'):
            if re.search(r"^@@@", each_body_text):
                n_tf_at +=1
        TF_AT = True if n_tf_at > 0 else False 
            
             
        if (auto_paragraphs_two == True) and (layout_idx == 3) and (TF_AT == False):
            body_t = body_text 
            body_t = remove_ptn_in_text(r"^<section:", body_t)
            body_t = remove_ptn_in_text(r"^<width:", body_t)
            body_t = remove_ptn_in_text(r"^<fig:", body_t)

            temp = format_text_in_two_columns(body_t, threshold=auto_paragraphs_threshold)
            #print(temp)
            body_21 = temp[0]
            body_22 = temp[1]
        else:
            body_text_list_at = body_text.split('@@@',1)
            body_21 = body_text_list_at[0]
            body_21 = remove_ptn_in_text(r"^<section:", body_21)
            body_21 = remove_ptn_in_text(r"^<width:", body_21)
            body_21 = remove_ptn_in_text(r"^<fig:", body_21)

            body_22 = body_text_list_at[1].strip() if len(body_text_list_at)>1 else ''

        body_total = body_21 +'\n' + body_22
        #print("*", len(body_11), len(body_21), len(body_22), len(body_total))
        title.text = slide_data['title']            
        
        # figure position 
        fig_pos_xpa=auto_figure_xp
        fig_pos_ypa=auto_figure_yp
        fig_height=auto_figure_height


        jj = 0 
        for each_text in body_text_list:
            jj +=1
         
            # 1. width 
            # 2. fig 
            # 3. 텍스트프레임에 add_paragrap()

            # 1 
            if re.search(r"<width:", each_text):
                #print(each_text)
                match1 = extract_width_height(each_text, height=default_textbox_height)
                #print(match1)
                if match1 :
                    #print(f"*match1 = {match1}")
                    body.width  = Inches(match1[0])
                    body.height = Inches(match1[1])
                    if layout_idx == 3:
                        body_right.width  = Inches(match1[0])
                        body_right.height = Inches(match1[1])
                continue
            # 2 section 
            section_show, section_color, section_width = extract_section_info(each_text, color="#eeeeee", width=0.5)
            if section_show:
                #elif re.search(r"^<section:True>", each_text):
                left = Inches(0)
                top = Inches(0)
                width = Inches(section_width)
                height = Inches(slide_height)
                text_box_top = slide.shapes.add_textbox(left, top, width, height)    
                # 텍스트 상자의 배경색 설정
                fill = text_box_top.fill
                fill.solid()  # 배경을 단색으로 설정
                fill.fore_color.rgb = hex_to_rgb_color(section_color)
                #print(f"w={body.width}, h={body.height}")
                continue 

            # 2 그림 지정 
            if re.search(r"<fig:", each_text):
                match = extract_figure_info(each_text)
                if match :
                    #print("*******************figure find")
                    file_fig = match[0]
                    fig_caption = match[4]
                    
                    if (auto_figure_position == True) and (auto_figure_vertical == True):
                        ic.enable()
                        #print("*auto positioning")
                        fig_pos_xpa=auto_figure_xp
                        fig_pos_ypa += auto_figure_delta

                        xp = fig_pos_xpa
                        yp = fig_pos_ypa
                        h = auto_figure_height 
                    elif (auto_figure_position == True) and (auto_figure_vertical == False):
                        ic.enable()
                        #print("*auto positioning")
                        fig_pos_xpa += auto_figure_delta
                        fig_pos_ypa  = auto_figure_yp 

                        xp = fig_pos_xpa
                        yp = fig_pos_ypa
                        h = auto_figure_height     
                    else:
                        
                        xp = match[1]
                        yp = match[2]
                        h = match[3]
                    #print(f"*** {xp}, {yp}, {h}")     

                    

                    add_figure_to_slide_one(slide=slide, 
                            dir_figure=dir_figure,
                            fig_name=file_fig,
                            color=COLORS['gray_silver'], 
                            pos_x=xp, 
                            pos_y=yp, 
                            height=h, 
                            caption=fig_caption, 
                            lw=0.1,                            
                            font_size_caption_title=font_size_caption_title)

                    

                    
                                   

                    # add_figure_to_slide(slide=slide, 
                    #             each_text=each_text, 
                    #             dir_figure=dir_figure,
                    #             color=COLORS['gray_silver'],
                    #             font_size_caption_title=font_size_caption_title)
                continue

            # 3 각라인 속성 지정 
            # elif each_text.strip() != '':
            #     # 패러그랩 추가하기, 패러그랩 레벨 속성 지정
            #     if layout_idx == 1:
            #         text_frame_add_paragraph_with_level(body_text_frame, each_text)           


        if layout_idx == 0 :
            format_text(title, font_name=font_name, font_size=80, color=COLORS['gray_dark'])
            title.left = Inches(1) # 왼쪽 여백 
            title.top = Inches(1)  # 위쪽 여백
            title.width = Inches(slide_width-2)   # 타이틀 박스 폭 조정 (기본보다 넓게 설정)
            title.height = Inches(1.5) # 타이틀 박스 높이 조정
            
            body = slide.placeholders[1]
            body.text = body_text
            format_text(body, font_size=30)

            m1 = 2
            #subtitle = body 
            body.left = Inches(1+m1)   # 왼쪽 여백
            body.top = Inches(5)      # 위쪽 여백
            body.width = Inches(slide_width-2-2*m1)   # 타이틀 박스 폭 조정 (기본보다 넓게 설정)
            body.height = Inches(1.5) # 타이틀 박스 높이 조정

        elif layout_idx in [1,3]: 
            # 1은 1단 
            # 3은 2단 
            # 제목, 본문(본문 스타일 지정)
            #title = slide.shapes.title
            title.left = Inches(0.2)   # 왼쪽 여백
            title.top = Inches(0.1)     
            title.height = Inches(1.2)
            title.width = Inches(slide_width-0.4)

            # [1] 타이틀 폰트 크기, 색 지정
            format_text(title, font_size=font_size_title, color=COLORS['gray_dark'])
            
        
            # [2] body에 있는 텍스트의 폰트를 조정
            # body 의 종류 : 1, 2  
            pos_top  = Inches(1.0)
            pos_left = Inches(1.0)
            pos_gap = Inches(0.5)
           
            func1 = partial(text_frame_add_paragraph_with_level,
                        font_name=font_name,
                        level_font_sizes=level_font_sizes,
                        space_before=space_before,
                        space_after=space_after, 
                        space_before1=space_before1,
                        space_after1=space_after1,
                        color_bold=color_bold
                         )
            match(layout_idx):
                case 1:
                    body = slide.shapes.placeholders[1]
                    body.left = pos_left
                    body.top = pos_top
                    body.width = default_textbox_width1
                    body.height= default_textbox_height1
                    body_text_frame = body.text_frame
                    body_text_frame.clear()
                    body.width = default_textbox_width1
                    for each_line in body_11.split('\n'):
                        func1(text_frame=body_text_frame, each_text=each_line)
                        # text_frame_add_paragraph_with_level(body_text_frame, 
                        #                                     each_text=each_line, 
                        #                                     font_name=font_name,
                        #                                     level_font_sizes=level_font_sizes,
                        #                                     space_before=space_before,
                        #                                     space_after=space_after, 
                        #                                     space_before1=space_before1,
                        #                                     space_after1=space_after1,
                        #                                     color_bold=color_bold )

                case 3:
                    body = slide.shapes.placeholders[1]
                    body.left = pos_left
                    body.top = pos_top
                    body.width = default_textbox_width_half
                    body.height= default_textbox_height1

                    body_right = slide.shapes.placeholders[2]
                    body_right.left = pos_left + default_textbox_width_half + pos_gap
                    body_right.top = pos_top
                    body_right.width = default_textbox_width_half
                    body_right.height= default_textbox_height1  

                    body_text_frame = body.text_frame
                    body_text_frame.clear()
                    body_text_frame_right= body_right.text_frame
                    body_text_frame_right.clear()

                    for each_line in body_21.split('\n'):
                        #body_text_frame.clear()
                        func1(text_frame=body_text_frame, each_text=each_line)
                        # text_frame_add_paragraph_with_level(body_text_frame, each_line, 
                        #                                     font_name=font_name,
                        #                                     level_font_sizes=level_font_sizes,
                        #                                     space_before=space_before,
                        #                                     space_after=space_after, color_bold=color_bold )
                    for each_line in body_22.split('\n'):
                        #body_text_frame_right.clear()
                        func1(text_frame=body_text_frame_right, each_text=each_line)
                        # text_frame_add_paragraph_with_level(body_text_frame_right,each_line,
                        #                                     font_name=font_name,
                        #                                     level_font_sizes=level_font_sizes,
                        #                                     space_before=space_before,
                        #                                     space_after=space_after, color_bold=color_bold
                        #                                     )
                    body_text_frame_right.vertical_anchor = MSO_ANCHOR.TOP
                    body_text_frame_right.margin_top = 0
                    body_text_frame_right.margin_bottom = 0
                    body_text_frame_right.margin_left = 0
                    body_text_frame_right.margin_right = 0


            body_text_frame.vertical_anchor = MSO_ANCHOR.TOP
            body_text_frame.margin_top = 0
            body_text_frame.margin_bottom = 0
            body_text_frame.margin_left = 0
            body_text_frame.margin_right = 0
            
           
    # 파일 저장
    name1 = datetime.today().strftime('%Y%m%d_%H%M%S')
    name2 = f"{prefix}_{name1}.pptx"
    prs.save(name2)
    print(f"*file saved : {name2}")



def format_text_in_two_columns(text, threshold=500):
    """
    입력된 텍스트가 특정 임계값을 넘으면 2단 형식으로 편집
    - 줄나눔 문자열(\n) 기준으로 분할
    """
    lines = text.split('\n')
    
    first_part = []
    second_part = []
    current_length = 0
    
    for line in lines:
        if current_length + len(line) <= threshold:
            first_part.append(line)
            current_length += len(line)
        else:
            second_part.append(line)
    
    #formatted_text = "\n".join(first_part) + "\n\n" + "\n".join(second_part) if second_part else "\n".join(first_part)
    #return formatted_text
    first_p = '\n'.join(first_part)
    second_p = '\n'.join(second_part)
    return [first_p,second_p]
    


def set_paragraphs(body_text_frame, font_size_level_0=16, 
                   font_size_1=14, font_size_2=12, 
                   font_size_3=10, font_size_4=8,
                   font_name='나눔고딕', font_size_default=8, space_before=3, space_after=3):
    # 레벨 기준으로 색, 폰트 사이즈 지정해 주기
    for paragraph in body_text_frame.paragraphs:
        #print(f"* paragraph={paragraph}")
        #paragraph.line_spacing = Pt(12)
        paragraph.space_before = Pt(space_before)  
        paragraph.space_after = Pt(space_after)
        paragraph.font.size = Pt(10) 
        for run in paragraph.runs:
            if paragraph.level == 0 : 
                run.font.size = Pt(font_size_level_0)  # 폰트 크기 설정
                run.font.bold = True  # 굵게 설정
                run.font.name = font_name # "나눔고딕"  
                run.font.color.rgb = gray_very_dark
            elif paragraph.level == 1 : 
                run.font.size = Pt(font_size_1)
                #run.font.bold = False
                run.font.name = font_name
                run.font.color.rgb = gray_dark
            elif paragraph.level == 2 : 
                run.font.size = Pt(font_size_2)
                run.font.name = font_name
                run.font.color.rgb = gray
            elif paragraph.level == 3: 
                run.font.size = Pt(font_size_3)
                run.font.color.rgb = gray 
            elif paragraph.level == 4: 
                run.font.size = Pt(font_size_4)
                run.font.color.rgb = gray 

            else:
                run.font.size = Pt(font_size_default) 
                run.font.color.rgb = gray_very_dark # 색상 설정 (파란색)

        paragraph.alignment = PP_ALIGN.LEFT  #  정렬

def add_figure_to_slide_one(slide, 
                            dir_figure:Path,
                            fig_name,
                            color, 
                            pos_x, pos_y, height, caption, 
                            lw,
                            font_size_caption_title):
    #print(">>> add_figure_to_slide_one()")

    img_path = dir_figure / fig_name  # 이미지 파일 경로
    # if img_path.is_file():
    #     print(f"img_path is file : {img_path}")
    # else:
    #     print(f"{img_path} was not file")
    img_path2 = str(img_path.resolve())

    left = Inches(pos_x)  
    top  = Inches(pos_y)   
    height = Inches(height)  # 이미지 너비 (높이는 비율에 맞게 자동 조정)
    try:
        pic = slide.shapes.add_picture(img_path2, left, top, height=height)

        # 그림의 라인 활성화 및 스타일 설정
        pic.line.color.rgb = COLORS['gray_light'] # 라인 색상
        pic.line.width = None # Inches(0.00)  # 라인 두께 0.02 
        pic.line.dash_style = None  # 라인 스타일 (None은 실선)

        # 이미지 아래에 캡션 추가
        caption_top = top + pic.height + Inches(0.0)  # 이미지 바로 아래 위치
        caption_left = left
        caption_width = pic.width
        caption_height = Inches(0.3)

        text_box_caption = slide.shapes.add_textbox(caption_left, caption_top, caption_width, caption_height)
        text_frame_caption = text_box_caption.text_frame

        # 캡션 텍스트 설정
        text_frame_caption.text = caption.strip() 

        # 텍스트 서식 조정
        text_frame_caption.paragraphs[0].font.size = Pt(font_size_caption_title)  # 글꼴 크기
        text_frame_caption.paragraphs[0].alignment = PP_ALIGN.LEFT #PP_ALIGN.CENTER  # 중앙 정렬
    except Exception as e:
        raise ValueError(f"Error in figure files : {img_path}")
    
def add_figure_to_slide(slide, each_text:str, dir_figure:Path,color=None, 
                        font_size_caption_title:int=8):
    if re.search(r"<fig:", each_text):
        #print(each_text)
        #figure_no += 1
        #print(f"*slide_no={slide_no}, figure_no : {figure_no}")
        #fig_name_list = re.split(r"<fig:|>", each_text2)
        #fig_name = ''.join(fig_name_list)
        #print(f"each_text2={each_text2}")
        matches = parse_figure_attributes_h(text=each_text)
        color = COLORS['gray_light'] if color is None else color 
        if matches:
            #print(f"*figure = {matches}")
            fig_name = matches[0]
            try:
                left1 = float(matches[1])
                top1=float(matches[2])
                height1=float(matches[3])
            except Exception as e:
                print(f"error : {e} in {fig_name}")
            
            if len(matches)==5:
                caption_text=matches[4].strip()
            else:
                caption_text=''
            if fig_name.endswith(('png','PNG','jpg','JPG')):
                #print(f">>> fg_name={fig_name}")

                img_path = dir_figure / fig_name  # 이미지 파일 경로
                # if img_path.is_file():
                #     print(f"img_path is file : {img_path}")
                # else:
                #     print(f"{img_path} was not file")
                img_path2 = str(img_path.resolve())

                left = Inches(left1)  # 왼쪽 여백
                top = Inches(top1)   # 위쪽 여백
                height = Inches(height1)  # 이미지 너비 (높이는 비율에 맞게 자동 조정)
                try:
                    pic = slide.shapes.add_picture(img_path2, left, top, height=height)

                    # 그림의 라인 활성화 및 스타일 설정
                    pic.line.color.rgb = COLORS['gray_light'] # 라인 색상
                    pic.line.width = None # Inches(0.00)  # 라인 두께 0.02 
                    pic.line.dash_style = None  # 라인 스타일 (None은 실선)

                    # 이미지 아래에 캡션 추가
                    caption_top = top + pic.height + Inches(0.0)  # 이미지 바로 아래 위치
                    caption_left = left
                    caption_width = pic.width
                    caption_height = Inches(0.3)

                    text_box_caption = slide.shapes.add_textbox(caption_left, caption_top, caption_width, caption_height)
                    text_frame_caption = text_box_caption.text_frame

                    # 캡션 텍스트 설정
                    text_frame_caption.text = caption_text.strip() 

                    # 텍스트 서식 조정
                    text_frame_caption.paragraphs[0].font.size = Pt(font_size_caption_title)  # 글꼴 크기
                    text_frame_caption.paragraphs[0].alignment = PP_ALIGN.LEFT #PP_ALIGN.CENTER  # 중앙 정렬
                except Exception as e:
                    raise ValueError(f"Error in figure files : {img_path}") 


def find_color_with_name(color_name:str="#ff0000"):
    
    match color_name.lower():
        case "dark_blue":
            color_bold = "#003366"
        case "dark_gray":
            color_bold = "#333333"
        case "dar_red":
            color_bold = "#8B0000"
        case "teal":
            color_bold = "#008080"
        case _ :
            if is_valid_hex_color(color_name):
                color_bold = color_name 
            else:
                color_bold = "#ff0000"
    return color_bold 

def is_valid_hex_color(color_code):
    pattern = r'^#[0-9A-Fa-f]{6}$'  # #뒤에 6자리 16진수 (0-9, A-F)
    return bool(re.match(pattern, color_code))

def hex_to_rgb_color(hex_code: str):
    """
    HEX 색상(#RRGGBB)을 RGBColor 객체로 변환하는 함수

    예제 입력: "#ff0000"
    반환: RGBColor(255, 0, 0)
    """
    # HEX 코드에서 '#' 제거 후 R, G, B 값 추출
    hex_code = hex_code.lstrip('#')
    
    if len(hex_code) != 6:
        raise ValueError("HEX 색상 코드는 '#RRGGBB' 형식이어야 합니다.")

    r, g, b = tuple(int(hex_code[i:i+2], 16) for i in (0, 2, 4))
    
    return RGBColor(r, g, b)


def text_frame_add_paragraph_with_level(text_frame, each_text: str = '', font_name='나눔고딕',
                                        level_font_sizes=None, 
                                        space_before=3, space_after=3,
                                        space_before1=1, space_after1=1,
                                        color_bold = "#ff0000",
                                        default_levl=1):
    """text frame에 문자열을 입력하여 레벨을 지정하고 들여쓰기를 적용하는 함수"""
    #print(">>> text_frame_add_paragraph_with_level()")
    p = text_frame.add_paragraph()
    #p.clear()
    # p.first_line_indent = 0  # 기본 들여쓰기 값 (첫 번째 줄)
    # p.left_indent = Pt(0)  # 기본 왼쪽 들여쓰기 값 (전체)

    each_text = each_text.strip()

    # bold color 지정해주기 
    color_bold = find_color_with_name(color_bold)
    

    # 레벨별 폰트 크기 및 들여쓰기 설정
    level_font_sizes_default = {
        0: Pt(24),  # 제목 수준
        1: Pt(20),
        2: Pt(18),
        3: Pt(16),
        4: Pt(14)   # 가장 낮은 수준
    }
    level_font_sizes = level_font_sizes if level_font_sizes is not None else level_font_sizes_default

    level_indentation = {
        0: Pt(0),  # 레벨 0: 기본 들여쓰기
        1: Pt(0),  # 레벨 1: 첫 번째 줄 + 왼쪽 20pt 들여쓰기
        2: Pt(0),  # 레벨 2: 첫 번째 줄 + 왼쪽 40pt 들여쓰기
        3: Pt(0),  # 레벨 3: 첫 번째 줄 + 왼쪽 60pt 들여쓰기
        4: Pt(0)   # 레벨 4: 첫 번째 줄 + 왼쪽 80pt 들여쓰기
    }

    # 계층(level) 설정
    level_patterns = [
        (r'^\*', 0),
        (r'^\-\-\-\-', 4),
        (r'^\-\-\-', 3),
        (r'^\-\-|^>>', 2), # type added 
        (r'^\-|^>', 1),    # type added 
        (r'^>', 3)
    ]

    # level setting 
    for pattern, level in level_patterns:
        match = re.match(pattern, each_text)
        if match:
            each_text = each_text[match.end():]  # 패턴 제거
            p.level = level
            break
        else:
            p.level = default_levl

        

    t1 = each_text.strip()
    if len(t1) > 0:
        # 1 pattenrn split 
        pattern = re.split(r"(\*\*.*?\*\*)", t1)  # ** 포함된 부분과 일반 텍스트를 분리

        # 들여쓰기 적용 (레벨별)
        level_indent = level_indentation.get(p.level, Pt(0))
        # if p.level != 0 : 
        #     print(f'p.level={p.level}, level indentation = {level_indent} ' )
        # else:
        #     print("*")
              

        #global crimson_red
        for part in pattern:
            run = p.add_run()
            set_font_style(run.font, font_name, level_font_sizes.get(p.level, Pt(18)))  # 레벨별 폰트 크기 적용

            if part.startswith("**") and part.endswith("**"):
               run.text = part[2:-2]  # ** 제거한 텍스트
               run.font.bold = True  # **로 감싸진 부분은 굵게
               color_bold2 = color_bold[1:]
               run.font.color.rgb = RGBColor(*tuple(int(color_bold2[i:i+2], 16) for i in (0, 2, 4))) # RGBColor(255, 0, 0) #crimson_red 
            else:
                run.text = part  # 일반 텍스트
            run.left_index = level_indent
            
            #run.text = part[2:-2] if part.startswith("**") and part.endswith("**") else part
        
        p.first_line_indent = level_indent  # 첫 번째 줄 들여쓰기
        p.left_indent = level_indent  # 전체 왼쪽 들여쓰기
        #print(f" left index = {level_indent}")
        match p.level:

            case 0:
                #print(f"p.level={p.level}")
                p.font.bold = True 
                p.bullet = None
                p.space_before = Pt(space_before)  
                p.space_after = Pt(space_after)
                #print(f"p.level={p.level}")
            case 1 : 
                p.space_before = Pt(space_before1)  
                p.space_after = Pt(space_after1)
            case 2 :
                p.space_before = Pt(space_before1)  
                p.space_after = Pt(space_after1)

            case _ :
                # p.bullet = 'o'
                p.space_before = Pt(3)  
                p.space_after = Pt(3)
            
    # if p.level != 0 : 
    #     print(f'- p.level={p.level}, level indentation = {level_indent} ' )
    # else:
    #     print("-")      
    return p

def format_text(title_shape, font_name='나눔고딕', font_size=20, color=COLORS["gray_dark"]):
    """타이틀 서식 지정"""
    for paragraph in title_shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = font_name
            run.font.size = Pt(font_size)
            run.font.bold = True
            run.font.color.rgb = color

def set_font_style(font, font_name, font_size):
    """공통 폰트 스타일 적용 함수 (레벨별 폰트 크기 적용)"""
    font.name = font_name
    font.size = font_size
    font.color.rgb = RGBColor(0, 0, 0)  # 검정색



def find_title_body(text:str, ptn=r'^#[\s\w]', symbol='//'):
    """
    # 타이틀
    그 아래는 body 
    """
    #textlist = re.split(ptn, text)
    #textlist = text.split('#')
    textlist = re.split(r'(?m)^#', text) # 각 라인 제일 처음 등장하는 # 기준으로 분리함 
    textlist1 = [each for each in textlist if each.strip() != '']
    dict_list = []
    for each in textlist:
        if each.strip() == '':
            continue
        temp = each.split('\n',1)
        if len(temp)==2:
            #print("***")
            title = temp[0]
            body = temp[1].strip()
        else:
            title = temp[0]
            body = f"* Error :{''.join(temp)}"
        
        #print(f"body={body}")
        body = re.sub(r'\t', ' ', body)
        body1 = body.split('\n')
        body2 = [b for b in body1 if b.strip() != '']
        body2 = '\n'.join(body1)
        #body = re.sub(r'\s+', ' ', body)
        #print(f"body={body}")
        body3 = remove_comment(body2, symbol)
        dict1 = {'title': title, 'body': body3}
        
        dict_list.append(dict1)
       
    return dict_list

def remove_comment(text, symbol):
    """
    주어진 텍스트에서 특정 부호로 시작하는 라인 삭제후 리턴함.
    Args:
        text (str): 입력 문자열
        symbol (str): 특정 부호 (예: '@', '*', '-')
    Returns:
        str: 주석 부문 삭제한 문자열
    """
    lines = text.splitlines()
    commented_lines = [ '' if line.lstrip().startswith(symbol) else line for line in lines]

    return "\n".join(commented_lines)



def read_text_file_for_pptx(filename):
    with open(filename, 'r', encoding='utf-8') as fn:
        text = fn.read()
        dict_list = find_title_body(text.strip())
        return dict_list



def find_figure(text=''):
    ptn1 = r"<fig:([\w\d\.\(\)]+),([\d\.]+),([\d\.]+),([\d\.]+)>" 
    ptn2 = r"<fig:([\w\d\.\(\)]+),([\d\.]+),([\d\.]+),([\d\.]+),([\s\w,\.\-\s\(\)\'\"]+)>" 
    matches  = re.findall(ptn1, text)
    matches2 = re.findall(ptn2, text)
   
    if matches:
        #print(matches[0])
        if len(matches[0]) == 4:
            return matches[0]
        else:   
            print(matches[0])
            raise ValueError('Error in Fig')
    elif matches2:
        #print(matches[0])
        if len(matches2[0]) == 5:
            return matches2[0]
        else:   
            #print(matches2[0])
            raise ValueError('Error in Fig')
    else:
        return []

def parse_figure_attributes(text, x=6.5,y=1, width=3, caption=''):
    text2 = [each for each in re.split(r'[:,<>]', text) if each.strip() != '']
    n = len(text2)
    match n :
        case 6:
            return text2[1:]
        case 5:
            return text2[1:] + [caption] 
        case 4:
            return text2[1:] + [width, caption]
        case 3:
            return text2[1:] + [y, width, caption]
        case 2:
            return text2[1:] + [x, y, width, caption]
        case _ :
            return None 

def parse_figure_attributes_h(text, x=6.5,y=1, height=3, caption=''):
    text2 = [each for each in re.split(r'[:,<>]', text) if each.strip() != '']
    n = len(text2)
    match n :
        case 6:
            return text2[1:]
        case 5:
            return text2[1:] + [caption] 
        case 4:
            return text2[1:] + [height, caption]
        case 3:
            return text2[1:] + [y, height, caption]
        case 2:
            return text2[1:] + [x, y, height, caption]
        case _ :
            return None 

def extract_figure_info(text: str):
    """
    주어진 문자열에서 <fig:> 패턴을 사용하여 파일명, 좌표1, 좌표2, 높이, 캡션을 추출하는 함수

    예제 입력:

    반환:
    ('time_schumpeterian.jpg', 4.0, 3.0, 2.8, 'Schumpeterian, ☞ 신슘페테리언(Neo-Schumpeterian) 이라고도 함')
    """
    text = text.strip() 
    pattern = r"<fig:\s*([\w\d\-_\.]+)\s*,\s*([\d\.]+)\s*,\s*([\d\.]+)\s*,\s*([\d\.]+)\s*,\s*(.*)\s*>"
    pattern2 = r"<fig:\s*([\w\d\-_\.]+)\s*,\s*([\d\.]+)\s*,\s*([\d\.]+)\s*,\s*([\d\.]+),?\s*>"

    match = re.search(pattern, text)
    if match:
        filename = match.group(1)
        x_coord = float(match.group(2))
        y_coord = float(match.group(3))
        height = float(match.group(4))
        caption = match.group(5)
        return filename, x_coord, y_coord, height, caption
    else:
        match = re.search(pattern2, text)
        if match:
            filename = match.group(1)
            x_coord = float(match.group(2))
            y_coord = float(match.group(3))
            height = float(match.group(4))
            caption = ''
            return filename, x_coord, y_coord, height,caption

    return None


def extract_section_info(text: str, color="#ff0000", width=0.5):
    """ section 정보 추출하는 함수 : color, width"""
    text = text.strip() 
    
    pattern_color = r"<section:(\w+),*\s*color=[\'\"]?(\w+)[\'\"]?,*width=([\d.]+)>"
    pattern =  r"<section:(\w+),*\s*(#\w{6})\s*,*([\d.]*)>"
    pattern2 = r"<section:(\w+),?\s*(#\w{6})\s*>"
    pattern3 =  r"<section:(\w+)>"

    match = re.search(pattern_color, text, flags=re.IGNORECASE)
    if match :
        show=match.group(1)
        color=match.group(2)
        color = COLORS_NAME.get(color, "#eeeeee")
        width=float(match.group(3))
    else:
        match = re.search(pattern, text, flags=re.IGNORECASE)
        show = False 
        if match:
            #print("match")
            show = match.group(1)
            color = match.group(2)
            width = float(match.group(3))
        else:
            #print("match2")
            match = re.search(pattern2, text, flags=re.IGNORECASE)
            if match :
                show = match.group(1)
                color = match.group(2)
            else:
                match = re.search(pattern3, text, flags=re.IGNORECASE)
                if match :
                    show = match.group(1)
    if show in ['true','True', 'TRUE', True]:
        show = True 
    else: 
        show = False 

    return show, color, width 


def extract_width_height(text='', ptn=None, width=8, height=5):
    ptn = r"<width:([\d\.]+),\s?height:([\w\.]+)>" if ptn is None else ptn 
    #print(f'a={a}')
    aa = re.findall(ptn, text)
    #print(aa)
    if aa:
        if len(aa[0]) == 2:
            width  = aa[0][0]
            height = aa[0][1]
        else:
            width = aa[0][0]
            height = height
        return float(width), float(height)
    else:
        return None 

def extract_slide_layout(text: str, default_index: int = 1) -> int:
    """슬라이드 레이아웃 번호 추출"""
    match = re.search(r"<slide layout index[:=](\d+)>", text)
    return int(match.group(1)) if match else default_index


def find_slide_layout_number(txt: str = '', default_index: int = 1) -> int:
    """
    Extracts the slide layout index from a given text.
    
    Layout indices:
    0 - Title (presentation title slide)
    1 - Title and Content
    2 - Section Header (default, sometimes called Segue)
    3 - Two Content (side by side bullet textboxes)
    4 - Comparison (additional title for each side-by-side content box)
    5 - Title Only
    6 - Blank
    7 - Content with Caption
    8 - Picture with Caption
    
    Example format in input text:
    "<slide layout index:3>"
    "<slide layout index=5>"
    
    :param txt: Input text containing the slide layout index
    :param default_index: Default index if no valid number is found
    :return: Extracted slide layout index as an integer
    """
    if not txt:  # 예외 처리
        return default_index

    for line in txt.split('\n'):
        cleaned_line = line.strip()
        if cleaned_line:
            match = re.match(r"^<slide layout index[:=](\d+)>", cleaned_line)
            if match:
                return int(match.group(1))  # 정수 변환 후 반환

    return default_index  # 기본값 반환


def clean_text(text: str, kind=1) -> str:
    """슬라이드 본문에서 특정 패턴 제거"""
    match kind:
        case 1 : 
            ptn =r"<.*?>"
        case 2|'layout_index':
            ptn =  r"<slide layout index[:=](\d+)>"
        case _ :
            return text.strip() 
    # clean pattern
    text2 = re.sub(ptn, '', text).strip()
    # clean double space 
    text3 = [each for each in text2.split('\n') if each.strip()!='']
    return '\n'.join(text3)

def remove_ptn_in_text( ptn: str = None, text: str='') -> str:
    """
    특정 패턴과 일치하는 행을 제거한 텍스트를 반환.

    :param text: 입력 텍스트
    :param ptn: 제거할 패턴 (기본값: "<slide layout index:x>" 형식)
    :return: 패턴이 제거된 텍스트
    """
    ptn = ptn or r"^<slide layout index[:=](\d+)>"  # 기본 패턴 설정

    text2 = '\n'.join(each for each in text.splitlines() if not re.match(ptn, each.strip()))
    return text2.strip()


def find_name(text):
    """find : name_kor, name_eng, birth_year, death_year
    
    [Example]

    text1="# 장바티스트 세(Jean-Baptiste Say, 1767-)"
    text2="# 장바티스트 세(Jean-Baptiste Say, 1767-1832)"

    find_name(text1)
    """
    pattern = r'\s*(?P<name_kor>.+?)\s*\((?P<name_eng>.+?),\s*(?P<birth_year>\d{4})-\)' 
    match  = re.search(pattern, text)
    dict1 = None 
    if match:
        name_kor = match.group('name_kor').strip()
        name_eng = match.group('name_eng').strip()
        birth_year = match.group('birth_year').strip()
        death_year = datetime.today().year
        dict1 = {'name_kor':name_kor, 'name_eng':name_eng, 'birth_year':birth_year, 'death_year':death_year}
        
    else:
        pattern = r'\s*(?P<name_kor>.+?)\s*\((?P<name_eng>.+?),\s*(?P<birth_year>\d{4})-\s*(?P<death_year>\d{4})\s*\)'
        match = re.search(pattern, text)
        if match :
            name_kor = match.group('name_kor').strip()
            name_eng = match.group('name_eng').strip()
            birth_year = match.group('birth_year').strip()
            death_year = match.group('death_year').strip()
            dict1 = {'name_kor':name_kor, 'name_eng':name_eng, 'birth_year':birth_year, 'death_year':death_year}
    return dict1
 
def read_files_for_collecting_names(dir_name=r'./text', pattern='*.txt'):
    """pattern 디렉토리 아래에 있는
    
    call : find_name()
    """
    names = []
    total_dict_list = []
    p = Path(dir_name)
    for each in p.glob(pattern):
        #print(each) 
        with open(each, 'r', encoding='utf-8') as fn:
            each_text = fn.read().split('*',1)
            each_text2 = each_text[-1].split('\n')[0]
            dict1 = find_name(each_text2)
            if dict1:
                dict1['start'] = datetime.strptime(str(dict1['birth_year']), "%Y")
                dict1['end'] = datetime.strptime(str(dict1['death_year']), "%Y")
                dict1['family_name'] = dict1['name_eng'].split(' ')[-1].strip()
                total_dict_list.append(dict1)
        

    df = DataFrame(total_dict_list)
    return df 

if __name__ == '__main__':

    print(__name__)

